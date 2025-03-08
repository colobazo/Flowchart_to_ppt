import cv2
import numpy as np
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

class FlowchartConverter:
    def __init__(self):
        self.image = None
        self.gray = None
        self.binary = None
        self.shapes = []
        self.connections = []
        self.texts = []
        self.shape_id_counter = 0  # Add counter for unique IDs

    def load_image(self, image_path):
        """Load and preprocess the image."""
        self.image = cv2.imread(image_path)
        if self.image is None:
            raise ValueError("Could not load image")
        
        # Convert to grayscale
        self.gray = cv2.cvtColor(self.image, cv2.COLOR_BGR2GRAY)
        
        # Apply thresholding
        _, self.binary = cv2.threshold(self.gray, 127, 255, cv2.THRESH_BINARY_INV)

    def detect_shapes(self):
        """Detect shapes in the image."""
        contours, _ = cv2.findContours(self.binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        
        for contour in contours:
            # Filter small contours
            if cv2.contourArea(contour) < 100:
                continue
                
            # Approximate the contour to detect shape
            epsilon = 0.04 * cv2.arcLength(contour, True)
            approx = cv2.approxPolyDP(contour, epsilon, True)
            
            # Get bounding box
            x, y, w, h = cv2.boundingRect(contour)
            
            # Determine shape type
            shape_type = self._classify_shape(approx, w/h)
            
            # Add shape with unique ID
            self.shapes.append({
                'id': self.shape_id_counter,  # Add unique ID
                'type': shape_type,
                'contour': contour,
                'position': (x, y),
                'size': (w, h)
            })
            self.shape_id_counter += 1

    def _classify_shape(self, approx, aspect_ratio):
        """Classify shape based on number of vertices and aspect ratio."""
        num_vertices = len(approx)
        
        if num_vertices == 4:
            if 0.95 <= aspect_ratio <= 1.05:
                return MSO_AUTO_SHAPE_TYPE.RECTANGLE
            else:
                return MSO_AUTO_SHAPE_TYPE.RECTANGLE
        elif num_vertices == 3:
            # Use diamond instead of triangle since triangle is not available
            return MSO_AUTO_SHAPE_TYPE.DIAMOND
        elif num_vertices > 6:
            return MSO_AUTO_SHAPE_TYPE.OVAL
        else:
            return MSO_AUTO_SHAPE_TYPE.RECTANGLE

    def extract_text(self):
        """Extract text from within shapes using OCR."""
        for shape in self.shapes:
            x, y, w, h = cv2.boundingRect(shape['contour'])
            roi = self.gray[y:y+h, x:x+w]
            
            # Create a mask for the shape
            mask = np.zeros_like(self.gray)
            cv2.drawContours(mask, [shape['contour']], -1, (255, 255, 255), -1)
            roi_mask = mask[y:y+h, x:x+w]
            
            # Apply mask to ROI
            roi_text = cv2.bitwise_and(roi, roi, mask=roi_mask[y:y+h, x:x+w])
            
            # Extract text using Tesseract
            text = pytesseract.image_to_string(roi_text).strip()
            shape['text'] = text

    def detect_connections(self):
        """Detect connections between shapes using line detection."""
        edges = cv2.Canny(self.gray, 50, 150)
        lines = cv2.HoughLinesP(edges, 1, np.pi/180, 50, minLineLength=30, maxLineGap=10)
        
        if lines is not None:
            for line in lines:
                x1, y1, x2, y2 = line[0]
                # Convert points to float tuples
                start_point = (float(x1), float(y1))
                end_point = (float(x2), float(y2))
                start_shape = self._find_connected_shape(start_point)
                end_shape = self._find_connected_shape(end_point)
                
                if start_shape and end_shape and start_shape != end_shape:
                    self.connections.append({
                        'start': start_shape,
                        'end': end_shape
                    })

    def _find_connected_shape(self, point):
        """Find which shape contains or is closest to the given point."""
        x, y = point
        point = (float(x), float(y))  # Convert to float tuple
        for shape in self.shapes:
            if cv2.pointPolygonTest(shape['contour'], point, False) >= 0:
                return shape
        return None

    def create_presentation(self, output_path):
        """Create PowerPoint presentation with detected shapes and connections."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        
        shape_objects = {}
        
        # Add shapes
        for shape in self.shapes:
            x, y = shape['position']
            w, h = shape['size']
            
            # Convert coordinates to inches (assuming standard 96 DPI)
            left = Inches(x / 96)
            top = Inches(y / 96)
            width = Inches(w / 96)
            height = Inches(h / 96)
            
            # Add shape to slide
            shape_obj = slide.shapes.add_shape(
                shape['type'],
                left, top, width, height
            )
            
            # Add text to shape
            if 'text' in shape and shape['text']:
                text_frame = shape_obj.text_frame
                text_frame.text = shape['text']
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
            shape_objects[shape['id']] = shape_obj  # Use ID as key
        
        # Add connections
        for connection in self.connections:
            start_shape = shape_objects[connection['start']['id']]  # Use ID to lookup shape
            end_shape = shape_objects[connection['end']['id']]  # Use ID to lookup shape
            
            connector = slide.shapes.add_connector(
                MSO_AUTO_SHAPE_TYPE.LINE_CONNECTOR_1,
                start_shape.left, start_shape.top,
                end_shape.left, end_shape.top
            )
            
            connector.begin_connect(start_shape, 0)
            connector.end_connect(end_shape, 0)
        
        prs.save(output_path)

def convert_flowchart(input_path, output_path):
    """Main function to convert flowchart image to PowerPoint."""
    converter = FlowchartConverter()
    converter.load_image(input_path)
    converter.detect_shapes()
    converter.extract_text()
    converter.detect_connections()
    converter.create_presentation(output_path)

if __name__ == "__main__":
    import sys
    
    if len(sys.argv) != 3:
        print("Usage: python flowchart_to_pptx.py <input_image> <output_pptx>")
        sys.exit(1)
        
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    
    try:
        convert_flowchart(input_path, output_path)
        print(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        print(f"Error: {str(e)}")
        sys.exit(1) 